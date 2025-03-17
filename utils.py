import json
from tqdm import tqdm
import fitz
from pathlib import Path
import base64
import os
import re
from nltk.translate.bleu_score import sentence_bleu, SmoothingFunction
from rouge_score import rouge_scorer
from typing import List, Dict
from multiprocessing import Pool, cpu_count
from PIL import Image
from pptx import Presentation

def extract_text_from_ppt(ppt_path: str) -> str:

    presentation = Presentation(ppt_path)

    all_content = {}
    for slide_number, slide in enumerate(presentation.slides):
        content = ""
        for shape in slide.shapes:

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    content += paragraph.text + "\n"
                content += "\n"
        all_content[str(slide_number)] = content

    return all_content

def encode_images_to_pil(img_dir):
    image_files = os.listdir(img_dir)
    image_files = sorted(image_files, key=lambda x: int(x.split("/")[-1].split('.')[0]))
    images = []
    
    for filename in image_files:
        images.append(Image.open(os.path.join(img_dir, filename)).convert("RGB"))
    
    return images

def cal_text_metrics(reference: str, hypothesis: str) -> Dict[str, float]:
    """
    Calculate BLEU and ROUGE scores between reference and hypothesis texts.
    
    Args:
        reference: Reference text
        hypothesis: Hypothesis text to compare against reference
        
    Returns:
        Dictionary containing BLEU and ROUGE scores
    """
    # Prepare texts for BLEU scoring
    ref_tokens = reference.lower().split()
    hyp_tokens = hypothesis.lower().split()
    
    # Calculate BLEU score with smoothing
    smooth = SmoothingFunction()
    bleu_score = sentence_bleu(
        [ref_tokens], 
        hyp_tokens,
        smoothing_function=smooth.method1
    )
    
    # Calculate ROUGE scores
    scorer = rouge_scorer.RougeScorer(['rouge1', 'rouge2', 'rougeL'], use_stemmer=True)
    rouge_scores = scorer.score(reference, hypothesis)
    
    return {
        'bleu': bleu_score,
        'rouge1': rouge_scores['rouge1'].fmeasure,
        'rouge2': rouge_scores['rouge2'].fmeasure, 
        'rougeL': rouge_scores['rougeL'].fmeasure
    }

def extract_json(text):
    """
    Extracts and parses JSON from a text string that may contain other content.
    
    This function attempts multiple strategies to extract valid JSON:
    1. First tries to find complete JSON structures using regex patterns
    2. Falls back to finding outermost { } or [ ] if regex fails
    3. Cleans the text to fix common JSON formatting issues
    4. Attempts multiple parsing methods with increasing aggressiveness
    
    Args:
        text: Input text that may contain JSON
        
    Returns:
        Parsed JSON as a dictionary or list
        
    Raises:
        ValueError: If JSON cannot be parsed after all attempts
    """
    import json
    import re
    
    # Clean up the text first to handle common issues
    # Remove markdown code block markers if present
    text = re.sub(r'```(?:json)?|```', '', text)
    
    # Strategy 1: Use regex to find JSON structures
    # Try to find a complete JSON object with balanced braces
    json_pattern = r'(\{(?:[^{}]|(?R))*\})'
    array_pattern = r'(\[(?:[^\[\]]|(?R))*\])'
    
    # Since Python's re doesn't support recursion (?R), we'll use a simpler approach
    # Try to match the outermost JSON object or array
    obj_match = re.search(r'\{.*\}', text, re.DOTALL)
    arr_match = re.search(r'\[.*\]', text, re.DOTALL)
    
    json_candidates = []
    if obj_match:
        json_candidates.append(obj_match.group(0))
    if arr_match:
        json_candidates.append(arr_match.group(0))
    
    # If we found potential JSON structures via regex
    for json_str in json_candidates:
        # Try to parse it directly first
        try:
            return json.loads(json_str)
        except json.JSONDecodeError:
            # If direct parsing fails, we'll continue to more aggressive methods
            pass
    
    # Strategy 2: Fall back to the original method if regex didn't work
    # Find leftmost { or [
    left_curly = text.find('{')
    left_bracket = text.find('[')
    
    # Determine which comes first (if both exist)
    if left_curly != -1 and left_bracket != -1:
        left_pos = min(left_curly, left_bracket)
    elif left_curly != -1:
        left_pos = left_curly
    elif left_bracket != -1:
        left_pos = left_bracket
    else:
        raise ValueError("No JSON structure found in text")
    
    # Find rightmost } or ]
    right_curly = text.rfind('}')
    right_bracket = text.rfind(']')
    
    # Determine which comes last (if both exist)
    if right_curly != -1 and right_bracket != -1:
        right_pos = max(right_curly, right_bracket)
    elif right_curly != -1:
        right_pos = right_curly
    elif right_bracket != -1:
        right_pos = right_bracket
    else:
        raise ValueError("No JSON structure found in text")
    
    # Extract the JSON substring
    json_str = text[left_pos:right_pos+1]
    
    # Strategy 3: Clean up common issues before parsing
    # Fix line breaks after commas which are common in LLM outputs
    json_str = re.sub(r',\s*\n\s*', ', ', json_str)
    
    # Fix missing quotes around keys
    json_str = re.sub(r'(\{|\,)\s*([a-zA-Z0-9_]+)\s*:', r'\1 "\2":', json_str)
    
    # Remove any trailing commas before closing brackets or braces
    json_str = re.sub(r',\s*(\}|\])', r'\1', json_str)
    
    # Strategy 4: Try multiple parsing methods with increasing aggressiveness
    # First attempt: standard json.loads
    try:
        return json.loads(json_str)
    except json.JSONDecodeError as e:
        print(f"Standard JSON parsing failed: {e}")
        
        # Second attempt: Try with ast.literal_eval for more forgiving parsing
        try:
            import ast
            return ast.literal_eval(json_str)
        except (SyntaxError, ValueError) as e:
            print(f"AST literal_eval failed: {e}")
            
            # Third attempt: Try with json5 if available (most lenient JSON parser)
            try:
                try:
                    import json5
                except ImportError:
                    # If json5 is not installed, try to install it
                    import subprocess
                    import sys
                    print("json5 module not found, attempting to install...")
                    subprocess.check_call([sys.executable, "-m", "pip", "install", "json5"])
                    import json5
                
                return json5.loads(json_str)
            except Exception as e:
                print(f"JSON5 parsing failed: {e}")
                
                # Final fallback: raise a more detailed error
                raise ValueError(f"Failed to parse JSON after multiple attempts. Original text: {text[:100]}...")

def count_words(text):
    chinese_characters = re.findall(r'[\u4e00-\u9fff]', text)
    english_words = re.findall(r'\b[a-zA-Z]+\b', text)
    
    chinese_char_count = len(chinese_characters)
    english_word_count = len(english_words)
    
    total_count = chinese_char_count + english_word_count
    
    return total_count


def encode_image_to_base64(img_path):
    try:
        # Get file extension and corresponding MIME type
        extension = os.path.splitext(img_path)[1].lower()
        mime_types = {
            '.png': 'image/png',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.gif': 'image/gif',
            '.bmp': 'image/bmp',
            '.webp': 'image/webp'
        }
        mime_type = mime_types.get(extension, 'image/jpeg')  # default to jpeg if unknown
        
        with open(img_path, 'rb') as f:
            image_data = f.read()
            base64_data = base64.b64encode(image_data).decode('utf-8')
            return f"data:{mime_type};base64,{base64_data}"
            
    except Exception as e:
        print(f"Error processing image {img_path}: {str(e)}")
        return None

def encode_images_to_base64(img_dir):
    image_files = os.listdir(img_dir)
    # Filter for files that match the pattern of "1.png", "2.png", etc.
    image_files = [f for f in image_files if re.match(r'^\d+\.png$', f)]
    image_files = sorted(image_files, key=lambda x: int(x.split("/")[-1].split('.')[0]))
    image_urls = []
    
    for filename in image_files:
        img_path = os.path.join(img_dir, filename)
        image_url = encode_image_to_base64(img_path)
        if image_url:
            image_urls.append(image_url)
    
    return image_urls

def convert_pdf_to_png(input_file, output_dir):
    os.makedirs(output_dir, exist_ok=True)

    # Open the PDF file
    doc = fitz.open(input_file)

    # Iterate through each page of the PDF
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)  # Load the current page
        pix = page.get_pixmap()  # Render page to an image
        output_path = os.path.join(output_dir, f'{page_num + 1}.png')  # Define output path for the image
        pix.save(output_path)  # Save the image

    print(f"Conversion completed. Images are saved in '{output_dir}'.")

def pptx_to_pdf(input_file: str, output_dir: str) -> bool:
    """Convert a PPTX file to PDF using LibreOffice in Docker.
    
    Args:
        input_file: Path to the PPTX file
        output_dir: Directory to save the PDF file
        
    Returns:
        bool: True if conversion successful, False otherwise
    """
    try:
        # Create output directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)
        
        # Convert paths to absolute paths
        input_file = os.path.abspath(input_file)
        input_dir = os.path.dirname(input_file)
        
        # Run LibreOffice conversion in Docker
        cmd = f'docker run --rm -v "{input_dir}:/data" -v "{output_dir}:/output" "libreoffice-converter" libreoffice --headless --convert-to pdf --outdir /output "{os.path.basename(input_file)}"'
        
        # Execute command and check return code
        result = os.system(cmd)
        if result != 0:
            print(f"Error converting {input_file} to PDF")
            return False
            
        return True
        
    except Exception as e:
        print(f"Error during PPTX to PDF conversion: {str(e)}")
        return False

def _parallel_wrapper(args):
    """
    Wrapper function for parallel processing that unpacks arguments.
    
    Args:
        args: Tuple of (func, arg) where:
            - func: The function to execute
            - arg: The argument to pass to the function
    """
    func, arg = args[0], args[1]
    try:
        return func(arg) if not isinstance(arg, tuple) else func(*arg)
    except Exception as e:
        print(f"Error processing with args {args}: {str(e)}")
        return None

def parallel_process(func, args_list, num_processes=None):
    """
    Run a function across multiple processes for parallel processing.
    
    Args:
        func: The function to run in parallel
        args_list: List of arguments, each containing the argument for one function call
        num_processes: Number of processes to use (defaults to CPU count)
        
    Returns:
        List of results from the function calls
    """
    if num_processes is None:
        num_processes = cpu_count()
    
    # Prepare arguments for each item
    process_args = [(func, args) for args in args_list]
    
    with Pool(processes=num_processes) as pool:
        results = list(tqdm(
            pool.imap(_parallel_wrapper, process_args),
            total=len(args_list),
            desc=f"Processing {len(args_list)} items with {num_processes} processes"
        ))
    
    return results