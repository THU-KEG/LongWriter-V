from pptx import Presentation
from inference.api.gpt import GPT_Interface
from tqdm import tqdm

def extract_text_from_ppt(ppt_path: str) -> str:

    presentation = Presentation(ppt_path)

    all_content = {}
    for slide_number, slide in enumerate(presentation.slides):
        content = ""
        for shape in slide.shapes:

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    content += paragraph.text + "\n"
                content += "\n"
        all_content[str(slide_number)] = content

    return all_content

def gen_script(ppt_path, images, progress_manager):
    system_prompt = [
        {"role": "system", "content": """This agent speaks Chinese. Lecture Script Writer's primary function is to analyze PowerPoint (PPT) slides based on user inputs and the texts extracted from those slides. It then generates a script for teachers to teach students about the content illustrated on the page, assuming the role of the teacher who also made the slides. The script is intended for the teacher to read out loud, directly engaging with the audience without referring to itself as an external entity. It focuses on educational content, suitable for classroom settings or self-study. It emphasizes clarity, accuracy, and engagement in explanations, avoiding overly technical jargon unless necessary. The agent is not allowed to ask the user any questions even if the provided information is insufficient or unclear, ensuring the responses have to be a script. The script for each slide is limited to no more than two sentences, leaving most of the details to be discussed when interacting with the student's questions. The scripts for each slide has to be consistant to the previouse slide and it is important to make sure the agent's generated return can be directly joined as a fluent and continued script without any further adjustment. The agent should also never assume what is one the next slide before processing it. It adopts a friendly and supportive tone, encouraging learning and curiosity."""}
    ]
    
    messages = []
    
    extracted_text = extract_text_from_ppt(ppt_path)

    progress_manager.update_progress("Script Generation", 0, "Generating scripts...")

    ctx_len = 3
    scripts = []
    for slide_number, content in tqdm(extracted_text.items(), desc="Generating scripts..."):
        msg = {"role": "user", "content": [{"type": "text", "text": content},
                                         {"type": "image_url", "image_url": {"url": images[int(slide_number)]}}]}
        
        messages.append(msg)
        
        res = GPT_Interface.call(model='gpt-4o-2024-05-13', messages=system_prompt + messages)

        scripts.append(res)
        
        progress_manager.update_progress("Script Generation", int(slide_number) / len(extracted_text) * 100, f"Generating script for slide {slide_number}...")
        
        messages.append({"role": "assistant", "content": res})
        
        if len(messages) > ctx_len:
            messages = messages[-ctx_len * 2:]
    
    progress_manager.update_progress("Script Generation", 100, "Scripts generated!")
    progress_manager.cleanup()

    return scripts